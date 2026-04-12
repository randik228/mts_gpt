"""
PPTX builder — generates PowerPoint presentations from structured data.

Usage:
    from core.pptx_builder import generate_pptx
    path = generate_pptx(title, slides)
"""
import json
import logging
import os
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

FILES_DIR = Path(os.environ.get("DATA_DIR", "/app/data")) / "files"
FILES_DIR.mkdir(parents=True, exist_ok=True)

# MTS brand colors
_MTS_RED = RGBColor(0xE3, 0x06, 0x11)
_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
_ACCENT = RGBColor(0x3A, 0x3A, 0x5C)


def _set_slide_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def generate_pptx(title: str, slides: list[dict]) -> str:
    """
    Generate a .pptx file from structured slide data.

    Args:
        title: Presentation title (used for filename and title slide).
        slides: List of dicts with keys:
            - title (str): Slide heading
            - content (str): Bullet points separated by newlines
            - notes (str, optional): Speaker notes

    Returns:
        Filename (relative to FILES_DIR) for download URL construction.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Title slide ---
    slide_layout = prs.slide_layouts[6]  # blank layout
    sl = prs.slides.add_slide(slide_layout)
    _set_slide_bg(sl, _DARK_BG)

    # Red accent bar at top
    from pptx.util import Emu
    bar = sl.shapes.add_shape(
        1, 0, 0, slide_width, Emu(int(slide_height * 0.02))
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _MTS_RED
    bar.line.fill.background()

    # Title text
    txBox = sl.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = _WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Подготовлено с помощью GPTHub AI"
    p2.font.size = Pt(16)
    p2.font.color.rgb = _LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT

    # --- Content slides ---
    for i, slide_data in enumerate(slides):
        sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _set_slide_bg(sl, _DARK_BG)

        # Red accent bar
        bar = sl.shapes.add_shape(
            1, 0, 0, slide_width, Emu(int(slide_height * 0.012))
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _MTS_RED
        bar.line.fill.background()

        # Slide number
        num_box = sl.shapes.add_textbox(
            Inches(12.2), Inches(6.8), Inches(1), Inches(0.5)
        )
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = str(i + 1)
        num_p.font.size = Pt(12)
        num_p.font.color.rgb = _LIGHT_GRAY
        num_p.alignment = PP_ALIGN.RIGHT

        # Slide title
        heading = slide_data.get("title", f"Слайд {i + 1}")
        title_box = sl.shapes.add_textbox(
            Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.0)
        )
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = heading
        tp.font.size = Pt(28)
        tp.font.bold = True
        tp.font.color.rgb = _WHITE
        tp.alignment = PP_ALIGN.LEFT

        # Content bullets
        content = slide_data.get("content", "")
        content_box = sl.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2)
        )
        ctf = content_box.text_frame
        ctf.word_wrap = True

        lines = [l.strip() for l in content.split("\n") if l.strip()]
        for j, line in enumerate(lines):
            # Remove leading bullet markers
            clean = line.lstrip("-*\u2022 ").strip()
            if not clean:
                continue
            if j == 0:
                p = ctf.paragraphs[0]
            else:
                p = ctf.add_paragraph()
            p.text = f"\u2022  {clean}"
            p.font.size = Pt(18)
            p.font.color.rgb = _WHITE
            p.space_after = Pt(8)

        # Speaker notes
        notes = slide_data.get("notes", "")
        if notes:
            sl.notes_slide.notes_text_frame.text = notes

    # Save
    filename = f"{uuid.uuid4().hex[:12]}.pptx"
    filepath = FILES_DIR / filename
    prs.save(str(filepath))
    logger.info("PPTX generated: %s (%d slides)", filepath, len(slides))
    return filename


def parse_presentation_json(raw: str) -> tuple[str, list[dict]]:
    """
    Parse model output into (title, slides).
    Expects JSON: {"title": "...", "slides": [{"title": "...", "content": "..."}]}
    Also handles markdown-fenced JSON blocks.
    """
    # Strip markdown fences
    text = raw.strip()
    if text.startswith("```"):
        lines = text.split("\n")
        # Remove first and last fence lines
        lines = [l for l in lines if not l.strip().startswith("```")]
        text = "\n".join(lines).strip()

    data = json.loads(text)
    title = data.get("title", "Презентация")
    slides = data.get("slides", [])
    # Normalize slides
    for s in slides:
        if "content" not in s:
            s["content"] = ""
        if "title" not in s:
            s["title"] = ""
    return title, slides
